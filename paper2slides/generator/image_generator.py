"""
Image Generator

Generate poster/slides images from ContentPlan.
Supports two modes:
- "api": Direct API calls to image generation service (default)
- "prompt": Export prompts for manual generation via web interface
"""
import os
import json
import base64
import time
import logging
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple
from openai import OpenAI
from concurrent.futures import ThreadPoolExecutor, as_completed
from PIL import Image, ImageDraw, ImageFont
import io

# Constants for placeholder images
PLACEHOLDER_BG_COLOR = '#F3F4F6'  # Light gray background
PLACEHOLDER_BORDER_COLOR = '#2563EB'  # Blue border
PLACEHOLDER_TITLE_COLOR = '#2563EB'  # Blue for slide number
PLACEHOLDER_TEXT_COLOR = '#1F2937'  # Dark gray text
PLACEHOLDER_WIDTH = 1920
PLACEHOLDER_HEIGHT = 1080
PLACEHOLDER_BORDER_WIDTH = 4

# Filename constants
GENERATED_IMAGE_FILENAME = 'generated.png'
SLIDE_DIR_TEMPLATE = 'slide_{:02d}_images'
SLIDE_PROMPT_TEMPLATE = 'slide_{:02d}_prompt.txt'
INSTRUCTIONS_FILENAME = 'INSTRUCTIONS.md'

# Alternative filenames to check when importing (fallbacks for generated.png)
IMPORT_ALTERNATIVE_FILENAMES = [
    'generated.jpg',
    'generated.jpeg',
    'slide.png',
    'slide.jpg',
    'output.png',
]

# Cross-platform font paths (tried in order)
FONT_PATHS_BOLD = [
    # Linux
    '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
    '/usr/share/fonts/TTF/DejaVuSans-Bold.ttf',
    # macOS
    '/System/Library/Fonts/Supplemental/Arial Bold.ttf',
    '/Library/Fonts/Arial Bold.ttf',
    # Windows
    'C:/Windows/Fonts/arialbd.ttf',
    'C:/Windows/Fonts/segoeui.ttf',
]
FONT_PATHS_REGULAR = [
    # Linux
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    '/usr/share/fonts/TTF/DejaVuSans.ttf',
    # macOS
    '/System/Library/Fonts/Supplemental/Arial.ttf',
    '/Library/Fonts/Arial.ttf',
    # Windows
    'C:/Windows/Fonts/arial.ttf',
    'C:/Windows/Fonts/segoeui.ttf',
]


def _load_font(font_paths: list, size: int):
    """Load font from first available path, fallback to default."""
    for path in font_paths:
        try:
            return ImageFont.truetype(path, size)
        except (OSError, IOError):
            continue
    return ImageFont.load_default()


from .config import GenerationInput
from .content_planner import ContentPlan, Section
from ..prompts.image_generation import (
    STYLE_PROCESS_PROMPT,
    FORMAT_POSTER,
    FORMAT_SLIDE,
    POSTER_STYLE_HINTS,
    SLIDE_STYLE_HINTS,
    SLIDE_LAYOUTS_ACADEMIC,
    SLIDE_LAYOUTS_DORAEMON,
    SLIDE_LAYOUTS_DEFAULT,
    SLIDE_COMMON_STYLE_RULES,
    POSTER_COMMON_STYLE_RULES,
    VISUALIZATION_HINTS,
    CONSISTENCY_HINT,
    SLIDE_FIGURE_HINT,
    POSTER_FIGURE_HINT,
)


@dataclass
class GeneratedImage:
    """Generated image result."""
    section_id: str
    image_data: bytes
    mime_type: str


@dataclass
class ProcessedStyle:
    """Processed custom style from LLM."""
    style_name: str       # e.g., "Cyberpunk sci-fi style with high-tech aesthetic"
    color_tone: str       # e.g., "dark background with neon accents"
    special_elements: str # e.g., "Characters appear as guides" or ""
    decorations: str      # e.g., "subtle grid pattern" or ""
    valid: bool
    error: Optional[str] = None


def process_custom_style(client: OpenAI, user_style: str, model: str = None) -> ProcessedStyle:
    """Process user's custom style request with LLM."""
    model = model or os.getenv("LLM_MODEL", "openai/gpt-4o-mini")
    
    try:
        response = client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": STYLE_PROCESS_PROMPT.format(user_style=user_style)}],
            response_format={"type": "json_object"},
        )
        result = json.loads(response.choices[0].message.content)
        return ProcessedStyle(
            style_name=result.get("style_name", ""),
            color_tone=result.get("color_tone", ""),
            special_elements=result.get("special_elements", ""),
            decorations=result.get("decorations", ""),
            valid=result.get("valid", False),
            error=result.get("error"),
        )
    except Exception as e:
        return ProcessedStyle(style_name="", color_tone="", special_elements="", decorations="", valid=False, error=str(e))


class ImageGenerator:
    """Generate poster/slides images from ContentPlan.

    Supports two modes:
    - "api": Direct API calls to image generation service (default)
    - "prompt": Export prompts for manual generation via web interface
    """

    def __init__(
        self,
        api_key: Optional[str] = None,
        base_url: Optional[str] = None,
        model: str = "google/gemini-3-pro-image-preview",
        mode: str = "api",
        prompt_output_dir: Optional[str] = None,
    ):
        self.mode = mode
        self.prompt_output_dir = Path(prompt_output_dir) if prompt_output_dir else None
        self.api_key = api_key or os.getenv("IMAGE_GEN_API_KEY", "")
        self.base_url = base_url or os.getenv("IMAGE_GEN_BASE_URL", "https://openrouter.ai/api/v1")
        self.model = model

        # Initialize image generation client (only in API mode)
        if self.mode == "api":
            self.client = OpenAI(api_key=self.api_key, base_url=self.base_url)
        else:
            self.client = None

        # Initialize LLM client for style processing (needed in both modes for custom styles)
        # Uses RAG_LLM_API_KEY which is typically an OpenAI key
        llm_api_key = os.getenv("RAG_LLM_API_KEY", "")
        llm_base_url = os.getenv("RAG_LLM_BASE_URL", "https://api.openai.com/v1")
        if llm_api_key:
            self.llm_client = OpenAI(api_key=llm_api_key, base_url=llm_base_url)
        else:
            self.llm_client = None

        # Track exported prompts for reference chain
        self._exported_prompts: List[dict] = []
        self._slide_counter = 0
    
    def generate(
        self,
        plan: ContentPlan,
        gen_input: GenerationInput,
        max_workers: int = 1,
        save_callback = None,
    ) -> List[GeneratedImage]:
        """
        Generate images from ContentPlan.
        
        Args:
            plan: ContentPlan from ContentPlanner
            gen_input: GenerationInput with config and origin
            max_workers: Maximum parallel workers for slides (3rd+ slides run in parallel)
            save_callback: Optional callback function(generated_image, index, total) called after each image
        
        Returns:
            List of GeneratedImage (1 for poster, N for slides)
        """
        figure_images = self._load_figure_images(plan, gen_input.origin.base_path)
        style_name = gen_input.config.style.value
        custom_style = gen_input.config.custom_style
        
        # Process custom style with LLM if needed
        processed_style = None
        if style_name == "custom" and custom_style:
            # Use llm_client for style processing (works in both API and prompt modes)
            style_client = self.llm_client or self.client
            if style_client is None:
                raise ValueError(
                    "Custom style requires an LLM client. "
                    "Set RAG_LLM_API_KEY environment variable for prompt export mode."
                )
            processed_style = process_custom_style(style_client, custom_style)
            if not processed_style.valid:
                raise ValueError(f"Invalid custom style: {processed_style.error}")
        
        all_sections_md = self._format_sections_markdown(plan)
        all_images = self._filter_images(plan.sections, figure_images)
        
        if plan.output_type == "poster":
            # Prompt export mode doesn't support posters (only slides workflow)
            if self.mode == "prompt":
                raise ValueError(
                    "Prompt export mode only supports slides output. "
                    "Use --output slides with --export-prompts."
                )
            result = self._generate_poster(style_name, processed_style, all_sections_md, all_images)
            if save_callback and result:
                save_callback(result[0], 0, 1)
            return result
        else:
            return self._generate_slides(plan, style_name, processed_style, all_sections_md, figure_images, max_workers, save_callback)
    
    def _generate_poster(self, style_name, processed_style: Optional[ProcessedStyle], sections_md, images) -> List[GeneratedImage]:
        """Generate 1 poster image."""
        prompt = self._build_poster_prompt(
            format_prefix=FORMAT_POSTER,
            style_name=style_name,
            processed_style=processed_style,
            sections_md=sections_md,
        )
        
        image_data, mime_type = self._call_model(prompt, images)
        return [GeneratedImage(section_id="poster", image_data=image_data, mime_type=mime_type)]
    
    def _generate_slides(self, plan, style_name, processed_style: Optional[ProcessedStyle], all_sections_md, figure_images, max_workers: int, save_callback=None) -> List[GeneratedImage]:
        """Generate N slide images.

        Dispatches to mode-specific implementation:
        - prompt mode: Sequential export for manual generation
        - api mode: First 2 sequential, rest parallel for efficiency
        """
        # Select layout rules based on style
        if style_name == "custom":
            layouts = SLIDE_LAYOUTS_DEFAULT
        elif style_name == "doraemon":
            layouts = SLIDE_LAYOUTS_DORAEMON
        else:
            layouts = SLIDE_LAYOUTS_ACADEMIC

        if self.mode == "prompt":
            return self._generate_slides_prompt_mode(
                plan, style_name, processed_style, all_sections_md, figure_images, layouts, save_callback
            )
        else:
            return self._generate_slides_api_mode(
                plan, style_name, processed_style, all_sections_md, figure_images, layouts, max_workers, save_callback
            )

    def _generate_slides_prompt_mode(
        self,
        plan,
        style_name,
        processed_style: Optional[ProcessedStyle],
        all_sections_md,
        figure_images,
        layouts,
        save_callback=None,
    ) -> List[GeneratedImage]:
        """Generate slides in prompt export mode (sequential, for manual generation)."""
        results = []
        total = len(plan.sections)

        for i in range(total):
            section = plan.sections[i]
            section_md = self._format_single_section_markdown(section, plan)
            layout_rule = layouts.get(section.section_type, layouts["content"])

            prompt = self._build_slide_prompt(
                style_name=style_name,
                processed_style=processed_style,
                sections_md=section_md,
                layout_rule=layout_rule,
                slide_info=f"Slide {i+1} of {total}",
                context_md=all_sections_md,
            )

            section_images = self._filter_images([section], figure_images)

            # Export prompt with reference chain instructions
            image_data, mime_type = self._export_prompt(
                prompt=prompt,
                reference_images=section_images,
                slide_num=i + 1,
                total_slides=total,
                section_title=section.title,
            )

            generated_img = GeneratedImage(section_id=section.id, image_data=image_data, mime_type=mime_type)
            results.append(generated_img)

            if save_callback:
                save_callback(generated_img, i, total)

        # Generate INSTRUCTIONS.md after all prompts are exported
        self._generate_instructions_md(total)

        return results

    def _generate_slides_api_mode(
        self,
        plan,
        style_name,
        processed_style: Optional[ProcessedStyle],
        all_sections_md,
        figure_images,
        layouts,
        max_workers: int,
        save_callback=None,
    ) -> List[GeneratedImage]:
        """Generate slides in API mode (first 2 sequential, rest parallel)."""
        results = []
        total = len(plan.sections)
        style_ref_image = None  # Store 2nd slide as reference for all subsequent slides

        # Generate first 2 slides sequentially (slide 1: no ref, slide 2: becomes ref)
        for i in range(min(2, total)):
            section = plan.sections[i]
            section_md = self._format_single_section_markdown(section, plan)
            layout_rule = layouts.get(section.section_type, layouts["content"])

            prompt = self._build_slide_prompt(
                style_name=style_name,
                processed_style=processed_style,
                sections_md=section_md,
                layout_rule=layout_rule,
                slide_info=f"Slide {i+1} of {total}",
                context_md=all_sections_md,
            )

            section_images = self._filter_images([section], figure_images)
            reference_images = []
            if style_ref_image:
                reference_images.append(style_ref_image)
            reference_images.extend(section_images)

            image_data, mime_type = self._call_model(prompt, reference_images)

            # Save 2nd slide (i=1) as style reference
            if i == 1:
                style_ref_image = {
                    "figure_id": "Reference Slide",
                    "caption": "STRICTLY MAINTAIN: same background color, same accent color, same font style, same chart/icon style. Keep visual consistency.",
                    "base64": base64.b64encode(image_data).decode("utf-8"),
                    "mime_type": mime_type,
                }

            generated_img = GeneratedImage(section_id=section.id, image_data=image_data, mime_type=mime_type)
            results.append(generated_img)

            # Save immediately if callback provided
            if save_callback:
                save_callback(generated_img, i, total)

        # Generate remaining slides in parallel (from 3rd onwards)
        if total > 2:
            results_dict = {}

            def generate_single(i, section):
                section_md = self._format_single_section_markdown(section, plan)
                layout_rule = layouts.get(section.section_type, layouts["content"])

                prompt = self._build_slide_prompt(
                    style_name=style_name,
                    processed_style=processed_style,
                    sections_md=section_md,
                    layout_rule=layout_rule,
                    slide_info=f"Slide {i+1} of {total}",
                    context_md=all_sections_md,
                )

                section_images = self._filter_images([section], figure_images)
                reference_images = [style_ref_image] if style_ref_image else []
                reference_images.extend(section_images)

                image_data, mime_type = self._call_model(prompt, reference_images)
                return i, GeneratedImage(section_id=section.id, image_data=image_data, mime_type=mime_type)

            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = {
                    executor.submit(generate_single, i, plan.sections[i]): i
                    for i in range(2, total)
                }

                for future in as_completed(futures):
                    idx, generated_img = future.result()
                    results_dict[idx] = generated_img

                    # Save immediately if callback provided
                    if save_callback:
                        save_callback(generated_img, idx, total)

            # Append in order
            for i in range(2, total):
                results.append(results_dict[i])

        return results

    def _generate_instructions_md(self, total_slides: int):
        """Generate INSTRUCTIONS.md file with workflow guide."""
        if not self.prompt_output_dir:
            return

        instructions = f"""# Manual Image Generation Workflow

## Overview
This directory contains {total_slides} slide prompts for manual generation via Nano Banana Pro Chat.

## Workflow

### Step 1: Generate Slides in Order (Important!)
Generate slides **in order** (1 → {total_slides}) to maintain visual consistency.

The reference chain works as follows:
- **Slide 1**: No reference needed (establishes base style)
- **Slide 2**: Use Slide 1 as reference (becomes THE style reference)
- **Slides 3-{total_slides}**: All use Slide 2 as reference

### Step 2: For Each Slide

1. Open the prompt file: `slide_XX_prompt.txt`
2. Read the **REFERENCE CHAIN INSTRUCTION** section
3. Upload any reference images from `slide_XX_images/` directory
4. Copy the **RAW PROMPT** section into Nano Banana Pro Chat
5. Generate the image
6. Download and save as: `slide_XX_images/{GENERATED_IMAGE_FILENAME}`

### Step 3: After All Slides Generated

Run the import command to create the final PPTX:
```bash
python -m paper2slides --import-images {self.prompt_output_dir}
```

## File Structure

```
prompts/
├── {INSTRUCTIONS_FILENAME} (this file)
├── slide_01_prompt.txt
├── slide_01_images/
│   ├── ref_00_*.png (reference images from paper)
│   └── {GENERATED_IMAGE_FILENAME} (YOU CREATE THIS)
├── slide_02_prompt.txt
├── slide_02_images/
│   └── ...
...
└── slide_{total_slides:02d}_prompt.txt
```

## Tips for Best Results

1. **Consistency**: Always upload Slide 2's generated image for slides 3+
2. **Style**: Tell Nano Banana to "maintain exact same style as reference"
3. **Colors**: If colors drift, explicitly mention the color palette
4. **Quality**: If a slide doesn't match, regenerate it before continuing

## Estimated Time
- Per slide: ~2-3 minutes
- Total: ~{total_slides * 2.5:.0f}-{total_slides * 3:.0f} minutes

## Troubleshooting

**Slide style doesn't match:**
- Regenerate with stronger reference instruction
- Upload both Slide 1 and Slide 2 as references

**Image quality issues:**
- Ask for "higher resolution" or "4K quality"
- Simplify the prompt if too complex

**Reference images missing:**
- Check the `slide_XX_images/` directory
- Re-run prompt export if files are missing
"""

        instructions_path = self.prompt_output_dir / INSTRUCTIONS_FILENAME
        with open(instructions_path, "w", encoding="utf-8") as f:
            f.write(instructions)

        logging.getLogger(__name__).info(f"  Generated: {INSTRUCTIONS_FILENAME}")
    
    def _format_custom_style_for_poster(self, ps: ProcessedStyle) -> str:
        """Format ProcessedStyle into style hints string for poster."""
        parts = [
            ps.style_name + ".",
            "English text only.",
            "Use ROUNDED sans-serif fonts for ALL text.",
            "Characters should react to or interact with the content, with appropriate poses/actions and sizes - not just decoration."
            f"LIMITED COLOR PALETTE (3-4 colors max): {ps.color_tone}.",
            POSTER_COMMON_STYLE_RULES,
        ]
        if ps.special_elements:
            parts.append(ps.special_elements + ".")
        return " ".join(parts)
    
    def _format_custom_style_for_slide(self, ps: ProcessedStyle) -> str:
        """Format ProcessedStyle into style hints string for slide."""
        parts = [
            ps.style_name + ".",
            "English text only.",
            "Use ROUNDED sans-serif fonts for ALL text.",
            "Characters should react to or interact with the content, with appropriate poses/actions and sizes - not just decoration.",
            f"LIMITED COLOR PALETTE (3-4 colors max): {ps.color_tone}.",
            SLIDE_COMMON_STYLE_RULES,
        ]
        if ps.special_elements:
            parts.append(ps.special_elements + ".")
        return " ".join(parts)
    
    def _build_poster_prompt(self, format_prefix, style_name, processed_style: Optional[ProcessedStyle], sections_md) -> str:
        """Build prompt for poster."""
        parts = [format_prefix]
        
        if style_name == "custom" and processed_style:
            parts.append(f"Style: {self._format_custom_style_for_poster(processed_style)}")
            if processed_style.decorations:
                parts.append(f"Decorations: {processed_style.decorations}")
        else:
            parts.append(POSTER_STYLE_HINTS.get(style_name, POSTER_STYLE_HINTS["academic"]))
        
        parts.append(VISUALIZATION_HINTS)
        parts.append(POSTER_FIGURE_HINT)
        parts.append(f"---\nContent:\n{sections_md}")
        
        return "\n\n".join(parts)
    
    def _build_slide_prompt(self, style_name, processed_style: Optional[ProcessedStyle], sections_md, layout_rule, slide_info, context_md) -> str:
        """Build prompt for slide with layout rules and consistency."""
        parts = [FORMAT_SLIDE]
        
        if style_name == "custom" and processed_style:
            parts.append(f"Style: {self._format_custom_style_for_slide(processed_style)}")
        else:
            parts.append(SLIDE_STYLE_HINTS.get(style_name, SLIDE_STYLE_HINTS["academic"]))
        
        # Add layout rule, then decorations if custom style
        parts.append(layout_rule)
        if style_name == "custom" and processed_style and processed_style.decorations:
            parts.append(f"Decorations: {processed_style.decorations}")
        
        parts.append(VISUALIZATION_HINTS)
        parts.append(CONSISTENCY_HINT)
        parts.append(SLIDE_FIGURE_HINT)
        
        parts.append(slide_info)
        parts.append(f"---\nFull presentation context:\n{context_md}")
        parts.append(f"---\nThis slide content:\n{sections_md}")
        
        return "\n\n".join(parts)
    
    def _format_sections_markdown(self, plan: ContentPlan) -> str:
        """Format all sections as markdown."""
        parts = []
        for section in plan.sections:
            parts.append(self._format_single_section_markdown(section, plan))
        return "\n\n---\n\n".join(parts)
    
    def _format_single_section_markdown(self, section: Section, plan: ContentPlan) -> str:
        """Format a single section as markdown."""
        lines = [f"## {section.title}", "", section.content]
        
        for ref in section.tables:
            table = plan.tables_index.get(ref.table_id)
            if table:
                focus_str = f" (focus: {ref.focus})" if ref.focus else ""
                lines.append("")
                lines.append(f"**{ref.table_id}**{focus_str}:")
                lines.append(ref.extract if ref.extract else table.html_content)
        
        for ref in section.figures:
            fig = plan.figures_index.get(ref.figure_id)
            if fig:
                focus_str = f" (focus: {ref.focus})" if ref.focus else ""
                caption = f": {fig.caption}" if fig.caption else ""
                lines.append("")
                lines.append(f"**{ref.figure_id}**{focus_str}{caption}")
                lines.append("[Image attached]")
        
        return "\n".join(lines)
    
    def _load_figure_images(self, plan: ContentPlan, base_path: str) -> List[dict]:
        """Load figure images as base64."""
        images = []
        mime_map = {
            ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".png": "image/png", ".webp": "image/webp", ".gif": "image/gif"
        }
        
        for fig_id, fig in plan.figures_index.items():
            if base_path:
                img_path = Path(base_path) / fig.image_path
            else:
                img_path = Path(fig.image_path)
            
            if not img_path.exists():
                continue
            
            mime_type = mime_map.get(img_path.suffix.lower(), "image/jpeg")
            
            try:
                with open(img_path, "rb") as f:
                    img_data = base64.b64encode(f.read()).decode("utf-8")
                images.append({
                    "figure_id": fig_id,
                    "caption": fig.caption,
                    "base64": img_data,
                    "mime_type": mime_type,
                })
            except Exception:
                continue
        
        return images
    
    def _filter_images(self, sections: List[Section], figure_images: List[dict]) -> List[dict]:
        """Filter images used in given sections."""
        used_ids = set()
        for section in sections:
            for ref in section.figures:
                used_ids.add(ref.figure_id)
        return [img for img in figure_images if img.get("figure_id") in used_ids]

    def _create_placeholder_image(self, slide_num: int, title: str = "") -> Tuple[bytes, str]:
        """Create a placeholder image for prompt export mode."""
        # Create 16:9 placeholder
        img = Image.new('RGB', (PLACEHOLDER_WIDTH, PLACEHOLDER_HEIGHT), color=PLACEHOLDER_BG_COLOR)
        draw = ImageDraw.Draw(img)

        # Draw border
        draw.rectangle(
            [10, 10, PLACEHOLDER_WIDTH - 10, PLACEHOLDER_HEIGHT - 10],
            outline=PLACEHOLDER_BORDER_COLOR,
            width=PLACEHOLDER_BORDER_WIDTH,
        )

        # Build slide directory and filename using constants
        slide_dir = SLIDE_DIR_TEMPLATE.format(slide_num)
        prompt_file = SLIDE_PROMPT_TEMPLATE.format(slide_num)

        # Draw centered text
        text_lines = [
            f"SLIDE {slide_num:02d}",
            "",
            "Placeholder Image",
            "",
            "Generate this slide manually using:",
            prompt_file,
            "",
            "Then place the generated image as:",
            f"{slide_dir}/{GENERATED_IMAGE_FILENAME}",
        ]
        if title:
            text_lines.insert(1, f"({title})")

        # Load fonts with cross-platform support
        font_large = _load_font(FONT_PATHS_BOLD, 48)
        font_medium = _load_font(FONT_PATHS_REGULAR, 32)

        y_position = 200
        for i, line in enumerate(text_lines):
            font = font_large if i == 0 else font_medium
            color = PLACEHOLDER_TITLE_COLOR if i == 0 else PLACEHOLDER_TEXT_COLOR
            bbox = draw.textbbox((0, 0), line, font=font)
            text_width = bbox[2] - bbox[0]
            x = (PLACEHOLDER_WIDTH - text_width) // 2
            draw.text((x, y_position), line, fill=color, font=font)
            y_position += 60 if i == 0 else 45

        # Convert to bytes
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        return buffer.getvalue(), 'image/png'

    def _save_reference_images(self, slide_dir: Path, reference_images: List[dict]) -> List[str]:
        """Save reference images to slide directory and return filenames."""
        saved_files = []
        for i, img in enumerate(reference_images):
            if not img.get("base64"):
                continue
            fig_id = img.get("figure_id", f"image_{i}")
            # Sanitize filename
            safe_name = "".join(c if c.isalnum() or c in "-_" else "_" for c in fig_id)
            ext = ".png" if "png" in img.get("mime_type", "png") else ".jpg"
            filename = f"ref_{i:02d}_{safe_name}{ext}"
            filepath = slide_dir / filename

            img_data = base64.b64decode(img["base64"])
            with open(filepath, "wb") as f:
                f.write(img_data)
            saved_files.append(filename)
        return saved_files

    def _build_reference_chain_instruction(self, slide_num: int, total_slides: int) -> str:
        """Build instruction for reference chain (style consistency)."""
        if slide_num == 1:
            return """
## REFERENCE CHAIN INSTRUCTION
This is SLIDE 1 - no reference image needed.
After generating this slide, save it - it will be used as reference for slide 2.
"""
        elif slide_num == 2:
            return """
## REFERENCE CHAIN INSTRUCTION
Upload the generated image from SLIDE 1 as a reference attachment.
Tell Nano Banana: "Use the same visual style, colors, and icon design as this reference image"

After generating this slide, save it - it will be the STYLE REFERENCE for ALL remaining slides.
"""
        else:
            return f"""
## REFERENCE CHAIN INSTRUCTION
Upload the generated image from SLIDE 2 as a reference attachment.
Tell Nano Banana: "STRICTLY MAINTAIN: same background color, same accent color, same font style, same chart/icon style. Keep visual consistency with this reference image."

This ensures consistent visual style across all {total_slides} slides.
"""

    def _export_prompt(
        self,
        prompt: str,
        reference_images: List[dict],
        slide_num: int,
        total_slides: int,
        section_title: str = "",
    ) -> Tuple[bytes, str]:
        """Export prompt and reference images for manual generation."""
        logger = logging.getLogger(__name__)

        if not self.prompt_output_dir:
            raise ValueError("prompt_output_dir must be set for prompt export mode")

        self.prompt_output_dir.mkdir(parents=True, exist_ok=True)

        # Create slide-specific directory
        slide_dir = self.prompt_output_dir / SLIDE_DIR_TEMPLATE.format(slide_num)
        slide_dir.mkdir(exist_ok=True)

        # Save reference images
        saved_refs = self._save_reference_images(slide_dir, reference_images)

        # Build reference chain instruction
        ref_chain_instruction = self._build_reference_chain_instruction(slide_num, total_slides)

        # Build slide directory name using constant
        slide_dir_name = SLIDE_DIR_TEMPLATE.format(slide_num)

        # Build complete prompt file
        prompt_content = f"""# Slide {slide_num:02d} of {total_slides}
{f'## {section_title}' if section_title else ''}

{ref_chain_instruction}

## REFERENCE IMAGES TO UPLOAD
Directory: {slide_dir_name}/
"""
        if saved_refs:
            for ref_file in saved_refs:
                prompt_content += f"- {ref_file}\n"
        else:
            prompt_content += "- (No reference images for this slide)\n"

        prompt_content += f"""
## RAW PROMPT FOR NANO BANANA
Copy everything below this line into Nano Banana Pro Chat:

---

{prompt}

---

## AFTER GENERATION
1. Download the generated image
2. Save as: {slide_dir_name}/{GENERATED_IMAGE_FILENAME}
3. Continue to the next slide prompt
"""

        # Save prompt file
        prompt_file = self.prompt_output_dir / SLIDE_PROMPT_TEMPLATE.format(slide_num)
        with open(prompt_file, "w", encoding="utf-8") as f:
            f.write(prompt_content)

        logger.info(f"  Exported prompt: {prompt_file.name}")

        # Track for instructions generation
        self._exported_prompts.append({
            "slide_num": slide_num,
            "title": section_title,
            "prompt_file": str(prompt_file),
            "images_dir": str(slide_dir),
            "reference_images": saved_refs,
        })

        # Return placeholder image
        return self._create_placeholder_image(slide_num, section_title)

    def _call_model(self, prompt: str, reference_images: List[dict]) -> tuple:
        """Call the image generation model with retry logic."""
        logger = logging.getLogger(__name__)
        content = [{"type": "text", "text": prompt}]
        
        # Add each image with figure_id and caption label
        for img in reference_images:
            if img.get("base64") and img.get("mime_type"):
                fig_id = img.get("figure_id", "Figure")
                caption = img.get("caption", "")
                label = f"[{fig_id}]: {caption}" if caption else f"[{fig_id}]"
                content.append({"type": "text", "text": label})
                content.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:{img['mime_type']};base64,{img['base64']}"}
                })
        
        # Retry logic for API calls
        max_retries = 3
        retry_delay = 2  # seconds
        
        for attempt in range(max_retries):
            try:
                logger.info(f"Calling image generation API (attempt {attempt + 1}/{max_retries})...")
                
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[{"role": "user", "content": content}],
                    extra_body={"modalities": ["image", "text"]}
                )
                
                # Check if response is valid
                if response is None:
                    error_msg = "API returned None response - possible rate limit or API error"
                    logger.warning(f"{error_msg} (attempt {attempt + 1}/{max_retries})")
                    if attempt < max_retries - 1:
                        time.sleep(retry_delay * (attempt + 1))
                        continue
                    raise RuntimeError(error_msg)
                
                if not hasattr(response, 'choices') or not response.choices:
                    error_msg = f"API response has no choices: {response}"
                    logger.warning(f"{error_msg} (attempt {attempt + 1}/{max_retries})")
                    if attempt < max_retries - 1:
                        time.sleep(retry_delay * (attempt + 1))
                        continue
                    raise RuntimeError(error_msg)
                
                message = response.choices[0].message
                if hasattr(message, 'images') and message.images:
                    image_url = message.images[0]['image_url']['url']
                    if image_url.startswith('data:'):
                        header, base64_data = image_url.split(',', 1)
                        mime_type = header.split(':')[1].split(';')[0]
                        logger.info("Image generation successful")
                        return base64.b64decode(base64_data), mime_type
                
                error_msg = "Image generation failed - no images in response"
                logger.warning(f"{error_msg} (attempt {attempt + 1}/{max_retries})")
                if attempt < max_retries - 1:
                    time.sleep(retry_delay * (attempt + 1))
                    continue
                raise RuntimeError(error_msg)
                
            except Exception as e:
                logger.error(f"Error in API call (attempt {attempt + 1}/{max_retries}): {str(e)}")
                if attempt < max_retries - 1:
                    time.sleep(retry_delay * (attempt + 1))
                    continue
                raise
        
        raise RuntimeError("Image generation failed after all retry attempts")


def save_images_as_pdf(images: List[GeneratedImage], output_path: str):
    """
    Save generated images as a single PDF file.
    
    Args:
        images: List of GeneratedImage from ImageGenerator.generate()
        output_path: Output PDF file path
    """
    from PIL import Image
    import io
    
    pdf_images = []
    
    for img in images:
        # Load image from bytes
        pil_img = Image.open(io.BytesIO(img.image_data))
        
        # Convert RGBA to RGB (PDF doesn't support alpha)
        if pil_img.mode == 'RGBA':
            pil_img = pil_img.convert('RGB')
        elif pil_img.mode != 'RGB':
            pil_img = pil_img.convert('RGB')
        
        pdf_images.append(pil_img)
    
    if pdf_images:
        # Save first image and append the rest
        pdf_images[0].save(
            output_path,
            save_all=True,
            append_images=pdf_images[1:] if len(pdf_images) > 1 else [],
            resolution=100.0,
        )
        print(f"PDF saved: {output_path}")


def save_images_as_pptx(images: List[GeneratedImage], output_path: str, title: str = "Generated Presentation"):
    """
    Save generated images as a PowerPoint presentation.

    Each image becomes a full-slide background image (16:9 aspect ratio).

    Args:
        images: List of GeneratedImage from ImageGenerator.generate()
        output_path: Output PPTX file path
        title: Presentation title (optional)
    """
    from pptx import Presentation
    from pptx.util import Inches

    logger = logging.getLogger(__name__)

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 width
    prs.slide_height = Inches(7.5)    # 16:9 height

    # Set presentation title in metadata
    if title:
        prs.core_properties.title = title

    # Add blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for img in images:
        # Add slide
        slide = prs.slides.add_slide(blank_layout)

        # Save image temporarily to add to slide
        img_stream = io.BytesIO(img.image_data)

        # Add image as full-slide background
        slide.shapes.add_picture(
            img_stream,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

    # Save presentation
    prs.save(output_path)
    logger.info(f"PPTX saved: {output_path}")


def import_generated_images(prompt_dir: str, output_path: str):
    """
    Import manually generated images from prompt export directory into PPTX.

    Looks for generated.png files in each slide_XX_images/ directory.
    Also imports speaker notes from checkpoint_plan.json if available.

    Args:
        prompt_dir: Directory containing slide_XX_images/ subdirectories
        output_path: Output PPTX file path

    Returns:
        List of missing slides (if any)
    """
    import json
    from pptx import Presentation
    from pptx.util import Inches

    logger = logging.getLogger(__name__)
    prompt_path = Path(prompt_dir)

    # Find all slide directories
    slide_dirs = sorted(prompt_path.glob("slide_*_images"))
    if not slide_dirs:
        raise ValueError(f"No slide directories found in {prompt_dir}")

    # Load speaker notes from checkpoint_plan.json if available
    # Try multiple locations: parent (timestamp dir) and grandparent (style dir)
    speaker_notes = {}
    checkpoint_path = None
    for candidate in [
        prompt_path.parent / "checkpoint_plan.json",           # prompts/../
        prompt_path.parent.parent / "checkpoint_plan.json",    # prompts/../../
    ]:
        if candidate.exists():
            checkpoint_path = candidate
            break

    if checkpoint_path:
        try:
            with open(checkpoint_path, "r", encoding="utf-8") as f:
                checkpoint_data = json.load(f)
            plan = checkpoint_data.get("plan", {})
            sections = plan.get("sections", [])
            for section in sections:
                slide_id = section.get("id", "")  # e.g., "slide_01"
                if slide_id.startswith("slide_"):
                    slide_num = int(slide_id.split("_")[1])
                    title = section.get("title", "")
                    content = section.get("content", "")

                    # Check for enhanced narrative notes first (from --enhance-speaker-notes)
                    if "speaker_notes_narrative" in section:
                        notes_text = section["speaker_notes_narrative"]
                    else:
                        # Check for rich speaker notes
                        notes_data = section.get("speaker_notes", {})
                        if notes_data and notes_data.get("talking_points"):
                            # Format rich speaker notes
                            notes_parts = [f"## {title}\n"]

                            # Add talking points
                            talking_points = notes_data.get("talking_points", [])
                            if talking_points:
                                notes_parts.append("### Key Points:")
                                for point in talking_points:
                                    notes_parts.append(f"• {point}")
                                notes_parts.append("")

                            # Add key terms to emphasize
                            key_terms = notes_data.get("key_terms", [])
                            if key_terms:
                                notes_parts.append(f"**Emphasize:** {', '.join(key_terms)}")
                                notes_parts.append("")

                            # Add transition
                            transition = notes_data.get("transition", "")
                            if transition:
                                notes_parts.append(f"**Transition:** {transition}")
                                notes_parts.append("")

                            # Add duration
                            duration = notes_data.get("duration_minutes", 2)
                            notes_parts.append(f"⏱️ ~{duration} minutes")

                            notes_text = "\n".join(notes_parts)
                        else:
                            # Fallback to basic format
                            notes_text = f"{title}\n\n{content}" if title else content

                    speaker_notes[slide_num] = notes_text
            logger.info(f"Loaded speaker notes for {len(speaker_notes)} slides")
        except Exception as e:
            logger.warning(f"Could not load speaker notes: {e}")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    missing_slides = []
    imported_count = 0

    for slide_dir in slide_dirs:
        # Extract slide number from directory name
        dir_name = slide_dir.name  # e.g., "slide_01_images"
        slide_num = int(dir_name.split("_")[1])

        # Look for generated image
        generated_img = slide_dir / GENERATED_IMAGE_FILENAME
        if not generated_img.exists():
            # Try alternative filenames
            for alt_name in IMPORT_ALTERNATIVE_FILENAMES:
                alt_path = slide_dir / alt_name
                if alt_path.exists():
                    generated_img = alt_path
                    break

        if not generated_img.exists():
            missing_slides.append(slide_num)
            logger.warning(f"Missing generated image for slide {slide_num}")
            continue

        # Add slide with image
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(generated_img),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Add speaker notes if available
        if slide_num in speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes[slide_num]

        imported_count += 1

    if imported_count == 0:
        error_msg = (
            f"No generated images found in {prompt_dir}. "
            f"Expected '{GENERATED_IMAGE_FILENAME}' in each slide_XX_images/ directory."
        )
        logger.error(error_msg)
        raise ValueError(error_msg)

    prs.save(output_path)
    logger.info(f"PPTX saved: {output_path} ({imported_count} slides)")

    return missing_slides
